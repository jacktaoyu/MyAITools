#!/usr/bin/env python3
"""Generate competition presentation PPT for BMS AUTOSAR Code Generator."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Theme colors
PRIMARY = RGBColor(0x1A, 0x5F, 0x7A)      # teal
ACCENT = RGBColor(0xFF, 0x6B, 0x35)       # orange
DARK = RGBColor(0x1E, 0x29, 0x33)         # dark blue-gray
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)        # light gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_shape(slide, text, left, top, width, height, font_size=44, bold=True, color=DARK):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return shape


def add_body_text(slide, text, left, top, width, height, font_size=18, color=DARK, line_spacing=1.3):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.line_spacing = line_spacing
    return shape


def add_bullet_list(slide, items, left, top, width, height, font_size=20, color=DARK):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(12)
    return shape


def add_section_header(slide, title, subtitle=None):
    # top banner
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = PRIMARY
    banner.line.fill.background()

    # title
    add_title_shape(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
                   font_size=32, color=WHITE)

    if subtitle:
        add_body_text(slide, subtitle, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
                     font_size=16, color=PRIMARY)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # === Slide 1: Title ===
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # background accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_title_shape(slide, "BMS AUTOSAR", Inches(1), Inches(2), Inches(11), Inches(1),
                   font_size=54, color=DARK)
    add_title_shape(slide, "智能代码生成器", Inches(1), Inches(2.9), Inches(11), Inches(1),
                   font_size=48, color=PRIMARY)
    add_body_text(slide, "基于 RAG 知识库 + MISRA 质量门 + ARXML 图谱的电池管理系统 AUTOSAR 代码生成工具",
                 Inches(1), Inches(4.1), Inches(11), Inches(1), font_size=22, color=DARK)
    add_body_text(slide, "参赛作品展示",
                 Inches(1), Inches(5.8), Inches(11), Inches(0.5), font_size=18, color=ACCENT)

    # === Slide 2: Background ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "背景与痛点", "Background & Pain Points")
    bullets = [
        "BMS（电池管理系统）是新能源汽车核心安全软件，AUTOSAR Classic 是行业主流架构",
        "传统开发依赖资深工程师手写 SWC、BSW、RTE、ARXML，周期长、成本高",
        "领域知识分散在规范文档、历史项目、ARXML 模型中，难以被 AI 直接利用",
        "生成代码常出现 MISRA C:2012 违规、接口不一致、命名不规范等问题",
        "现有通用代码生成工具缺乏 AUTOSAR/BMS 领域语义理解能力"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 3: Solution Overview ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "方案概览", "Solution Overview")
    bullets = [
        "面向 BMS AUTOSAR 的端到端智能代码生成 VS Code 插件",
        "融合 RAG 知识库检索、领域 query 扩展、LLM 重排序，提升检索准确率",
        "内置 MISRA C:2012 静态检查 + 自动修复闭环，确保代码合规",
        "ARXML 知识图谱可视化，辅助理解模型结构与接口关系",
        "支持用户自定义模板，内置多种 BMS 组件类型（CSC、Controller、Balancer 等）"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 4: Architecture ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "系统架构", "System Architecture")

    boxes = [
        ("用户交互层", 0.6, 2.2, ["VS Code 插件", "Wizard 向导", "知识库管理器", "质量报告面板", "ARXML 图谱视图"]),
        ("任务编排层", 4.7, 2.2, ["BmsAutosarGenerateHandler", "进度事件总线", "gRPC 服务", "模板渲染引擎"]),
        ("智能检索层", 8.8, 2.2, ["Embedding 语义检索", "BM25 词法检索", "Query 扩展", "ARXML 图谱增强", "LLM 重排序"]),
        ("知识/质量层", 0.6, 4.8, ["BMS 知识库", "用户模板", "MISRA 检查器", "质量报告存储", "自动修复器"]),
        ("生成输出层", 8.8, 4.8, [".h / .c 源码", "ARXML 描述文件", "统一 diff 预览", "批量应用修复"]),
    ]

    for title, x, y, items in boxes:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.6), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = PRIMARY
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            p.space_after = Pt(4)

    # === Slide 5: Core Feature - Code Generation ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "核心功能 ①：智能代码生成", "Intelligent Code Generation")
    bullets = [
        "支持 12 种组件类型：SWC、BSW Module、RTE Interface、ARXML Descriptor、Service 等",
        "BMS 专用组件：CSC（Cell Supervision Circuit / AFE）、Controller、Balancer、Thermal Manager、Charger、Diagnosis",
        "一键生成 .h / .c 源码 + ARXML 描述文件 + Types 定义",
        "支持批量配置（YAML/JSON）批量生成多个组件蓝图",
        "生成过程实时流式展示进度：preparing → retrieving_knowledge → generating → complete"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 6: Core Feature - RAG ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "核心功能 ②：RAG 知识库增强检索", "RAG Knowledge Base Retrieval")
    bullets = [
        "Hybrid 检索：Embedding 余弦相似度 + BM25 词法分融合，可调的 hybrid weight",
        "AUTOSAR 领域 Query 扩展：自动扩展 CSC/AFE、SOC/StateOfCharge、BMS/BatteryManagementSystem 等同义词",
        "元数据预过滤：按 tag（如 swc/cell/thermal）和 source file 精准过滤",
        "ARXML 图谱增强：从源文件构建 SWC/Port/Interface 图谱，给拓扑邻近条目加权",
        "LLM-as-reranker：二阶段 LLM 打分重排，进一步提升检索相关性"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 7: Core Feature - Quality ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "核心功能 ③：MISRA 质量门", "MISRA C:2012 Quality Gates")
    bullets = [
        "内置轻量级 MISRA C:2012 静态检查器，覆盖 12 条常见生成代码违规",
        "检测项包括：stdlib/stdio 禁用、八进制常量、未初始化变量、多返回点、函数声明缺失等",
        "生成代码后自动运行质量门，生成结构化质量报告",
        "Webview 质量报告面板按 error/warning/info 分级展示，支持一键定位源码",
        "支持单文件与批量自动修复，生成统一 diff 预览后用户确认再写盘"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 8: Core Feature - Knowledge Graph ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "核心功能 ④：ARXML 知识图谱", "ARXML Knowledge Graph")
    bullets = [
        "解析 ARXML 文件，提取 SWC、Port、Interface、Data Type、Runnable 等实体",
        "构建节点与边的关系网络，支持多跳邻居查询",
        "检索阶段利用图谱邻近度对相关知识条目加权，提升领域相关性",
        "Webview 交互式力导向图谱可视化，直观展示系统架构",
        "支持选择具体 ARXML 文件进行图谱构建"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 9: Core Feature - Templates ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "核心功能 ⑤：用户自定义模板", "User-Defined Templates")
    bullets = [
        "内置默认模板覆盖常见组件类型，保证开箱即用",
        "支持 workspace/global 两级用户模板存储（templates.json）",
        "模板优先级：workspace > global > built-in，用户可覆盖内置生成模式",
        "Wizard 中组件类型下拉框自动合并内置类型与用户模板",
        "提供模板管理对话框，支持创建、列出、删除自定义模板"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 10: Technical Highlights ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "技术亮点", "Technical Highlights")
    bullets = [
        "端到端 VS Code 插件形态，直接服务工程师日常开发流程",
        "gRPC + Protobuf 实现前后端高效通信，支持流式生成进度推送",
        "原子写入、embedding 缓存、并发控制等工程化细节保证稳定性",
        "模块化架构：检索、生成、检查、修复、图谱、模板各模块可独立演进",
        "RAG 参数可调 UI：Top K、Hybrid Weight、Score Threshold、LLM Reranker 实时配置"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 11: Validation ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "验证与测试", "Validation & Testing")
    bullets = [
        "TypeScript 类型检查：npm run check-types ✅",
        "代码规范：npm run lint ✅",
        "单元测试：npm run test:unit ✅ 1647 项通过",
        "插件构建：npm run package:vsix ✅ 生成可安装 .vsix 文件（9.1 MB）",
        "覆盖模块：Query Expander、Reranker、Semantic Retrieval Filters、Template Storage、Auto-Fixer 等"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 12: Demo / Results ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "成果展示", "Results")
    bullets = [
        "已实现完整的 BMS AUTOSAR 代码生成 VS Code 扩展",
        "支持从需求/组件类型到源码+ARXML 的一站式生成",
        "RAG 检索准确率通过领域扩展、图谱增强、LLM 重排序持续提升",
        "生成代码经过 MISRA 检查与自动修复闭环，质量可控",
        "构建产物可直接安装到 VS Code 进行演示和使用"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 13: Innovation ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "创新点", "Innovation Points")
    bullets = [
        "领域专用 RAG：针对 AUTOSAR/BMS 缩写与概念做 query 扩展，解决词汇鸿沟",
        "图谱增强检索：首次将 ARXML 知识图谱融入代码生成的上下文检索",
        "LLM 二阶段重排：在 embedding+BM25 后引入 LLM-as-reranker，兼顾效率与精度",
        "生成-检查-修复闭环：MISRA 检查与 LLM 自动修复形成完整质量闭环",
        "模板即配置：用户无需修改代码即可通过模板覆盖生成规则"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 14: Future Work ===
    slide = prs.slides.add_slide(slide_layout)
    add_section_header(slide, "未来展望", "Future Work")
    bullets = [
        "引入更多 AUTOSAR 标准规范（如 AUTOSAR CP R22-11）作为结构化知识",
        "增加生成代码的单元测试与仿真验证能力",
        "支持从自然语言需求自动推导端口、runnable 和数据类型",
        "建立 BMS 代码生成 benchmark，量化 RAG 与生成质量提升",
        "探索多 Agent 协作：架构设计 Agent + 代码生成 Agent + 审查 Agent"
    ]
    add_bullet_list(slide, bullets, Inches(0.7), Inches(2.2), Inches(12), Inches(5), font_size=22)

    # === Slide 15: Thank You ===
    slide = prs.slides.add_slide(slide_layout)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    add_title_shape(slide, "感谢聆听", Inches(0), Inches(2.8), Inches(13.333), Inches(1),
                   font_size=60, color=WHITE)
    add_body_text(slide, "BMS AUTOSAR 智能代码生成器",
                 Inches(0), Inches(4.1), Inches(13.333), Inches(0.6), font_size=28, color=WHITE)
    add_body_text(slide, "欢迎提问与交流",
                 Inches(0), Inches(5.0), Inches(13.333), Inches(0.5), font_size=20, color=ACCENT)

    return prs


if __name__ == "__main__":
    prs = create_presentation()
    prs.save("BMS_AUTOSAR_Code_Generator_Presentation.pptx")
    print("PPT generated: BMS_AUTOSAR_Code_Generator_Presentation.pptx")
